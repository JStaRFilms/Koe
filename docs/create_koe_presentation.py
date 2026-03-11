from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "Koe_Application_Overview.pptx"
LOGO = ROOT / "src" / "assets" / "icons" / "logo.png"

NAVY = RGBColor(13, 27, 42)
INK = RGBColor(23, 37, 53)
TEAL = RGBColor(39, 211, 195)
CYAN = RGBColor(129, 236, 236)
AMBER = RGBColor(255, 183, 3)
CORAL = RGBColor(243, 114, 82)
SAND = RGBColor(246, 242, 235)
WHITE = RGBColor(255, 255, 255)
SLATE = RGBColor(103, 116, 138)
PALE = RGBColor(230, 236, 242)
MINT = RGBColor(224, 251, 247)

TITLE_FONT = "Georgia"
BODY_FONT = "Aptos"


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def textbox(slide, x, y, w, h, text, size=18, color=INK, bold=False,
            font=BODY_FONT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def bullet_box(slide, x, y, w, h, lines, color=INK, size=18, accent=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for index, line in enumerate(lines):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.level = 0
        p.bullet = True
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = line
        run.font.name = BODY_FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
        if accent and index == 0:
            run.font.bold = True
    return box


def card(slide, x, y, w, h, fill, radius=True, line=None):
    shape_type = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def title_block(slide, title, subtitle=None, dark=False):
    color = WHITE if dark else INK
    sub_color = CYAN if dark else SLATE
    textbox(slide, 0.75, 0.45, 8.7, 0.75, title, size=30, color=color, bold=True, font=TITLE_FONT)
    if subtitle:
        textbox(slide, 0.78, 1.12, 8.7, 0.45, subtitle, size=12.5, color=sub_color)


def stat_card(slide, x, y, w, h, value, label, dark=False):
    bg = RGBColor(22, 37, 58) if dark else WHITE
    fg = WHITE if dark else INK
    sub = CYAN if dark else SLATE
    card(slide, x, y, w, h, bg, line=TEAL if dark else PALE)
    textbox(slide, x + 0.18, y + 0.18, w - 0.36, 0.55, value, size=22, color=fg, bold=True, font=TITLE_FONT)
    textbox(slide, x + 0.18, y + 0.78, w - 0.36, 0.35, label, size=10.5, color=sub)


def feature_card(slide, x, y, w, h, eyebrow, title, body):
    card(slide, x, y, w, h, WHITE, line=PALE)
    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.18), Inches(y + 0.18), Inches(0.42), Inches(0.42))
    badge.fill.solid()
    badge.fill.fore_color.rgb = TEAL
    badge.line.color.rgb = TEAL
    textbox(slide, x + 0.72, y + 0.18, w - 0.9, 0.22, eyebrow.upper(), size=8.5, color=SLATE, bold=True)
    textbox(slide, x + 0.18, y + 0.55, w - 0.36, 0.35, title, size=15, color=INK, bold=True)
    textbox(slide, x + 0.18, y + 0.98, w - 0.36, h - 1.1, body, size=10.5, color=SLATE)


def process_step(slide, x, y, w, h, number, title, body):
    card(slide, x, y, w, h, WHITE, line=PALE)
    bubble = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.18), Inches(y + 0.18), Inches(0.48), Inches(0.48))
    bubble.fill.solid()
    bubble.fill.fore_color.rgb = AMBER
    bubble.line.color.rgb = AMBER
    textbox(slide, x + 0.31, y + 0.24, 0.2, 0.2, str(number), size=11, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, x + 0.18, y + 0.8, w - 0.36, 0.35, title, size=14.5, color=INK, bold=True)
    textbox(slide, x + 0.18, y + 1.15, w - 0.36, h - 1.3, body, size=10.3, color=SLATE)


def arrow(slide, x, y, w):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(0.24))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.color.rgb = TEAL


def add_logo(slide, x, y, w):
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(x), Inches(y), width=Inches(w))


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.author = "OpenAI Codex"
    prs.core_properties.title = "Koe Application Overview"
    prs.core_properties.subject = "Application presentation for Koe desktop dictation app"

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(8.8), Inches(-0.3), Inches(5), Inches(8))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(18, 45, 67)
    accent.line.color.rgb = RGBColor(18, 45, 67)
    orb = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.1), Inches(0.75), Inches(3.2), Inches(3.2))
    orb.fill.solid()
    orb.fill.fore_color.rgb = RGBColor(25, 65, 92)
    orb.line.color.rgb = RGBColor(25, 65, 92)
    add_logo(slide, 9.85, 1.4, 1.75)
    textbox(slide, 0.8, 0.8, 5.8, 0.5, "APPLICATION OVERVIEW", size=11, color=CYAN, bold=True)
    textbox(slide, 0.8, 1.35, 6.8, 1.2, "Koe", size=31, color=WHITE, bold=True, font=TITLE_FONT)
    textbox(slide, 0.8, 2.18, 7.2, 1.0, "Lightning-fast, privacy-first voice dictation for Windows.", size=23, color=WHITE)
    textbox(slide, 0.8, 3.15, 6.7, 1.0, "Press a hotkey, speak naturally, and get polished text pasted into any focused app.", size=13.5, color=PALE)
    stat_card(slide, 0.8, 5.2, 1.8, 1.25, "Free", "BYOK with Groq", dark=True)
    stat_card(slide, 2.8, 5.2, 2.2, 1.25, "216x", "realtime Whisper speed", dark=True)
    stat_card(slide, 5.25, 5.2, 2.05, 1.25, "v1.0.7", "Electron desktop release", dark=True)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, SAND)
    title_block(slide, "Why Koe exists", "A practical alternative to subscription dictation tools.")
    textbox(slide, 0.8, 1.5, 3.3, 1.1, "The problem", size=18, color=CORAL, bold=True, font=TITLE_FONT)
    textbox(slide, 0.8, 1.95, 3.5, 1.7,
            "Most voice dictation tools either charge monthly, depend on the OS default experience, or skip the AI cleanup that makes dictated text usable.",
            size=14, color=INK)
    textbox(slide, 0.8, 4.15, 3.5, 1.2, "Koe fixes that with a local-first capture flow, Groq Whisper for transcription, and Llama cleanup for readable output.", size=14, color=SLATE)
    compare = [
        (4.55, "WhisperFlow", "Paid SaaS", "Strong UX\nMonthly cost"),
        (6.75, "OS Dictation", "Bundled", "Low accuracy\nMinimal formatting"),
        (8.95, "Koe", "Free BYOK", "High accuracy\nAI enhancement"),
    ]
    for x, name, price, detail in compare:
        fill = WHITE if name != "Koe" else MINT
        line = PALE if name != "Koe" else TEAL
        card(slide, x, 1.75, 1.95, 3.95, fill, line=line)
        textbox(slide, x + 0.18, 2.02, 1.55, 0.35, name, size=16, color=INK, bold=True, font=TITLE_FONT)
        textbox(slide, x + 0.18, 2.55, 1.4, 0.3, price, size=10.5, color=CORAL if name != "Koe" else TEAL, bold=True)
        for idx, line_text in enumerate(detail.split("\n")):
            textbox(slide, x + 0.18, 3.2 + idx * 0.48, 1.55, 0.35, line_text, size=13.5, color=SLATE)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    title_block(slide, "From hotkey to polished text", "A short pipeline designed for speed and low friction.", dark=True)
    steps = [
        ("Trigger", "Global shortcut starts capture from any app."),
        ("Detect", "Local Silero VAD keeps listening offline and waits for real speech."),
        ("Transcribe", "Groq Whisper converts audio to text with near-instant turnaround."),
        ("Enhance", "Llama cleans filler words, punctuation, and formatting."),
        ("Deliver", "Koe copies or auto-pastes the result where the cursor already is."),
    ]
    step_x = [0.65, 3.1, 5.55, 8.0, 10.45]
    for idx, ((title, body), x) in enumerate(zip(steps, step_x), start=1):
        process_step(slide, x, 2.15, 2.15, 2.6, idx, title, body)
        if idx < len(steps):
            arrow(slide, x + 1.98, 3.28, 0.48)
    note = card(slide, 3.45, 5.45, 6.4, 0.95, RGBColor(22, 37, 58), line=RGBColor(40, 77, 104))
    textbox(slide, 3.7, 5.72, 5.9, 0.3, "Privacy guardrail: audio only leaves the device after speech is detected.", size=12.5, color=CYAN, bold=True)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, SAND)
    title_block(slide, "Core product features", "The current release focuses on the daily dictation loop.")
    features = [
        ("Input", "Global hotkey", "Starts or stops dictation from anywhere with Ctrl+Shift+Space."),
        ("Audio", "Local VAD", "ONNX-powered speech detection runs on-device before any upload."),
        ("Speed", "Instant transcription", "Groq Whisper handles chunks at very high realtime speed."),
        ("Output", "AI cleanup", "Llama turns rough dictation into clean, readable copy."),
        ("Workflow", "Auto-paste", "Results can be typed directly into the active text field."),
        ("Recall", "History and usage", "Recent transcriptions and free-tier usage stay visible in-app."),
    ]
    coords = [(0.8, 1.65), (4.45, 1.65), (8.1, 1.65), (0.8, 4.15), (4.45, 4.15), (8.1, 4.15)]
    for (eyebrow, title, body), (x, y) in zip(features, coords):
        feature_card(slide, x, y, 3.0, 1.95, eyebrow, title, body)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_block(slide, "Architecture at a glance", "Electron splits system integration from the recording UI.")
    card(slide, 0.8, 1.65, 5.85, 4.9, RGBColor(244, 247, 250), line=PALE)
    card(slide, 6.95, 1.65, 5.55, 4.9, RGBColor(248, 251, 253), line=PALE)
    textbox(slide, 1.1, 1.95, 2.5, 0.35, "Renderer process", size=18, color=INK, bold=True, font=TITLE_FONT)
    textbox(slide, 7.25, 1.95, 2.6, 0.35, "Main process", size=18, color=INK, bold=True, font=TITLE_FONT)
    renderer_blocks = [
        (1.1, 2.55, 5.2, 0.8, TEAL, "Floating pill UI", "Always-on-top status, live feedback, minimal distraction."),
        (1.1, 3.55, 5.2, 0.8, CYAN, "Audio pipeline", "Web Audio API, AudioWorklet, VAD, and WAV encoding."),
        (1.1, 4.55, 5.2, 0.8, AMBER, "Settings and history", "Language, API key, enhancement mode, usage readout."),
    ]
    main_blocks = [
        (7.25, 2.55, 4.95, 0.8, TEAL, "System tray + shortcuts", "Runs quietly in the tray and listens for global commands."),
        (7.25, 3.55, 4.95, 0.8, CYAN, "Groq service layer", "Whisper transcription, Llama enhancement, rate limiting."),
        (7.25, 4.55, 4.95, 0.8, AMBER, "Clipboard and paste bridge", "Returns cleaned text back to the active app."),
    ]
    for x, y, w, h, fill, head, body in renderer_blocks + main_blocks:
        shape = card(slide, x, y, w, h, fill, line=fill)
        shape.fill.transparency = 0.18
        textbox(slide, x + 0.18, y + 0.1, w - 0.36, 0.23, head, size=12.5, color=INK, bold=True)
        textbox(slide, x + 0.18, y + 0.38, w - 0.36, 0.28, body, size=10.2, color=INK)
    connector = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(6.35), Inches(3.55), Inches(0.42), Inches(0.42))
    connector.fill.solid()
    connector.fill.fore_color.rgb = CORAL
    connector.line.color.rgb = CORAL
    textbox(slide, 6.2, 4.1, 0.8, 0.2, "IPC", size=10, color=CORAL, bold=True, align=PP_ALIGN.CENTER)
    chips = ["Electron", "Vite", "Whisper", "Llama 3", "vad-web", "electron-store"]
    chip_x = 0.95
    for label in chips:
        width = 0.8 + len(label) * 0.08
        shape = card(slide, chip_x, 6.75, width, 0.38, RGBColor(244, 247, 250), radius=False, line=PALE)
        textbox(slide, chip_x + 0.12, 6.84, width - 0.24, 0.18, label, size=9.2, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
        chip_x += width + 0.14

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    title_block(slide, "Privacy and economics", "The product stays useful because the operating model stays lightweight.", dark=True)
    card(slide, 0.85, 1.7, 5.25, 4.95, RGBColor(21, 38, 58), line=RGBColor(35, 67, 93))
    textbox(slide, 1.15, 2.0, 3.0, 0.35, "Privacy-first decisions", size=18, color=WHITE, bold=True, font=TITLE_FONT)
    bullet_box(
        slide, 1.15, 2.55, 4.55, 3.2,
        [
            "Voice activity detection runs locally with ONNX WebAssembly.",
            "Speech chunks are uploaded only when actual speaking is detected.",
            "No long-term audio storage is required for the core workflow.",
            "The user's Groq API key is stored locally with electron-store.",
        ],
        color=PALE, size=12.5
    )
    stat_card(slide, 6.55, 1.95, 2.05, 1.35, "20", "requests per minute", dark=True)
    stat_card(slide, 8.85, 1.95, 2.05, 1.35, "2,000", "requests per day", dark=True)
    stat_card(slide, 11.0, 1.95, 1.5, 1.35, "28.8k", "seconds of audio/day", dark=True)
    card(slide, 6.55, 3.7, 5.95, 2.1, RGBColor(21, 38, 58), line=RGBColor(35, 67, 93))
    textbox(slide, 6.85, 4.0, 5.2, 0.45, "Why this matters", size=18, color=WHITE, bold=True, font=TITLE_FONT)
    textbox(slide, 6.85, 4.48, 5.1, 1.0,
            "Koe avoids the usual SaaS tax by letting the user bring a free Groq key. That keeps infrastructure cost near zero while preserving a premium dictation experience.",
            size=13, color=PALE)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, SAND)
    title_block(slide, "Roadmap", "The foundation is in place; the next gains are workflow depth and platform reach.")
    columns = [
        (0.85, "Now", TEAL, ["Global hotkey toggle", "Local VAD and chunking", "Groq Whisper + Llama cleanup", "History and usage dashboard"]),
        (4.55, "Next", AMBER, ["Custom AI prompts", "Shortcut customization", "Export history as .txt or .md", "Mac support"]),
        (8.25, "Later", CORAL, ["Snippet library", "App-specific tone profiles", "Cloud sync", "Team collaboration"]),
    ]
    for x, label, color, items in columns:
        shape = card(slide, x, 1.85, 3.15, 4.85, WHITE, line=PALE)
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(1.85), Inches(3.15), Inches(0.28))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.color.rgb = color
        textbox(slide, x + 0.2, 2.25, 1.2, 0.3, label, size=18, color=INK, bold=True, font=TITLE_FONT)
        bullet_box(slide, x + 0.2, 2.8, 2.7, 3.4, items, color=SLATE, size=11.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_logo(slide, 0.95, 0.8, 1.0)
    textbox(slide, 2.1, 0.95, 4.5, 0.35, "KOE", size=18, color=CYAN, bold=True)
    textbox(slide, 0.95, 2.05, 8.0, 0.9, "A fast, focused dictation app built around one job: turning speech into usable text with as little friction as possible.", size=24, color=WHITE, font=TITLE_FONT, bold=True)
    textbox(slide, 0.95, 3.35, 7.0, 0.8, "The product combines desktop-level ergonomics, practical privacy constraints, and free-tier economics into a tool that can compete with paid alternatives for everyday use.", size=14, color=PALE)
    quote = card(slide, 8.55, 1.75, 3.75, 3.2, RGBColor(21, 38, 58), line=RGBColor(35, 67, 93))
    textbox(slide, 8.9, 2.2, 3.05, 0.45, "Positioning", size=16, color=CYAN, bold=True)
    textbox(slide, 8.9, 2.8, 2.95, 1.4, "\"WhisperFlow-quality convenience without the recurring subscription.\"", size=20, color=WHITE, font=TITLE_FONT)
    textbox(slide, 0.95, 6.45, 4.8, 0.25, "Repository version captured in this deck: 1.0.7", size=10.5, color=CYAN)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
